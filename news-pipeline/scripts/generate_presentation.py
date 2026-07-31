"""
Génère la présentation PowerPoint du projet (UTF-8, français).
Exécution : python scripts/generate_presentation.py
Sortie : docs/Presentation_Architecture_Donnees.pptx
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "docs" / "Presentation_Architecture_Donnees.pptx"

TEAM = "Zaidi Hamza — Elmer — Abdellah Mamdouh"


def add_title_slide(prs, title, subtitle):
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    if slide.placeholders[1].has_text_frame:
        slide.placeholders[1].text = subtitle


def add_bullet_slide(prs, title, bullets):
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    body = slide.placeholders[1].text_frame
    body.clear()
    for i, line in enumerate(bullets):
        if i == 0:
            p = body.paragraphs[0]
        else:
            p = body.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(20)


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_title_slide(
        prs,
        "Architecture de données",
        f"Collecte et analyse des actualités (web scraping)\n\nÉquipe : {TEAM}",
    )

    add_bullet_slide(
        prs,
        "Contexte",
        [
            "Volume d’articles médias exploitable pour tendances, thèmes, veille.",
            "Objectif : plateforme Big Data distribuée.",
            "Ingestion, stockage, transformation, qualité, gouvernance, visualisation.",
        ],
    )

    add_bullet_slide(
        prs,
        "Sources & collecte",
        [
            "Python + BeautifulSoup + requests.",
            "Sites Maroc + international (config sources_config.py).",
            "Champs : titre, auteur, date, catégorie, contenu, source, URL, pays…",
        ],
    )

    add_bullet_slide(
        prs,
        "Ingestion batch + streaming",
        [
            "Batch : DAG Airflow hourly + scraper SCRAPER_ONCE.",
            "Streaming : Kafka (raw-articles) → consumer → MinIO silver.",
            "Scraper : Kafka + MinIO bronze pour chaque article.",
        ],
    )

    add_bullet_slide(
        prs,
        "Data lake & médaillon",
        [
            "MinIO (S3-compatible) : buckets bronze / silver / gold.",
            "Bronze : JSON brut ; Silver : nettoyage + langue.",
            "Gold : tables analytiques PostgreSQL + résumé JSON dans MinIO.",
        ],
    )

    add_bullet_slide(
        prs,
        "Entrepôt & analyses",
        [
            "PostgreSQL (news_dw) : articles, par jour, source, thème, pays.",
            "Mots-clés pour tendances ; vues SQL pour Grafana.",
        ],
    )

    add_bullet_slide(
        prs,
        "Orchestration Airflow",
        [
            "DAG news_pipeline : scrape → bronze→silver → silver→gold → DQ.",
            "UI http://localhost:8080 — admin / admin (démo).",
            "Image Docker custom : dépendances pré-installées (démarrage fiable).",
        ],
    )

    add_bullet_slide(
        prs,
        "Qualité, gouvernance, visualisation",
        [
            "DQ : complétude, validité, cohérence → table dq_checks.",
            "data_catalog : inventaire des jeux de données.",
            "Grafana : tendances, volumes, mots-clés, pays, thèmes.",
        ],
    )

    add_bullet_slide(
        prs,
        "Si Airflow ne s’ouvre pas",
        [
            "docker compose ps / logs airflow — attendre 2–5 min après démarrage.",
            "Conflit port 8080 → mapper 8081:8080.",
            "Docker : augmenter la RAM ; tester http://127.0.0.1:8080.",
        ],
    )

    add_bullet_slide(
        prs,
        "Conclusion",
        [
            "Architecture alignée sur le cahier des charges.",
            "Pistes : fake news (ML), séparation DB Airflow / DW.",
            f"Merci — {TEAM}",
        ],
    )

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print("Fichier créé :", OUT)


if __name__ == "__main__":
    main()
